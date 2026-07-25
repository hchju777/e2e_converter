import unittest

from pptx import Presentation
from pptx.util import Emu

from src.dashboard import report_style as style
from src.dashboard.generate_pptx import (
    EMU_PER_IN,
    KOREAN_LINE_HEIGHT,
    MIN_TABLE_FONT_SIZE,
    _add_data_slide,
    _cell_spans,
    _fit_table_font,
    _table_height_for,
    _wrapped_line_count,
)
from src.dashboard.report_data import ReportPage, TableGrid


def report_page(rows: int, with_chart: bool, cells=None, merges=None) -> ReportPage:
    return ReportPage(
        section_id="test",
        banner=None,
        page_num="1",
        title="test",
        insight_lines=[],
        unit_note=None,
        table=TableGrid(rows, 2, cells or {}, merges or [], 1, set()),
        chart_png=b"png" if with_chart else None,
        chart_data=None,
        footnote=None,
    )


class ReportLayoutTests(unittest.TestCase):
    def test_chart_and_table_never_allocate_more_than_available_height(self):
        available = Emu(int(5 * EMU_PER_IN))
        table_height = _table_height_for(report_page(30, True), available, True)

        self.assertLessEqual(table_height, available * 0.48)

    def test_table_only_page_can_use_all_available_height(self):
        available = Emu(int(3 * EMU_PER_IN))
        table_height = _table_height_for(report_page(30, False), available, False)

        self.assertEqual(table_height, available)

    def test_insight_box_has_nonzero_height_and_does_not_overlap_unit_note(self):
        page = report_page(4, False)
        page.insight_lines = ["첫 번째 인사이트", "두 번째 인사이트"]
        page.unit_note = "[Unit: %]"
        prs = Presentation()
        _add_data_slide(prs, page)
        slide = prs.slides[0]
        insight = next(shape for shape in slide.shapes if shape.has_text_frame and "첫 번째" in shape.text)
        unit = next(shape for shape in slide.shapes if shape.has_text_frame and "[Unit" in shape.text)

        self.assertGreater(insight.height, 0)
        self.assertLessEqual(insight.top + insight.height, unit.top)


class TableFittingTests(unittest.TestCase):
    """표가 배정된 높이를 넘지 않아야 각주를 덮지 않는다."""

    COL_WIDTHS = [Emu(int(1.5 * EMU_PER_IN)), Emu(int(1.5 * EMU_PER_IN))]

    def _fit(self, rows, height_in, cells=None, merges=None):
        table = report_page(rows, False, cells, merges).table
        height = Emu(int(height_in * EMU_PER_IN))
        font, v_margin, row_lines = _fit_table_font(
            table, self.COL_WIDTHS, _cell_spans(table), height
        )
        total = sum(
            lines * font * KOREAN_LINE_HEIGHT / 72 * EMU_PER_IN + 2 * int(v_margin)
            for lines in row_lines
        )
        return font, row_lines, total, int(height)

    def test_long_table_still_fits_the_allotted_height(self):
        font, _, total, height = self._fit(rows=32, height_in=5.65)

        self.assertLessEqual(total, height)
        self.assertLessEqual(font, style.TABLE_FONT_SIZE)

    def test_shrinks_font_when_rows_do_not_fit(self):
        roomy, _, _, _ = self._fit(rows=10, height_in=5.0)
        cramped, _, total, height = self._fit(rows=40, height_in=4.0)

        self.assertLess(cramped, roomy)
        self.assertGreaterEqual(cramped, MIN_TABLE_FONT_SIZE)
        self.assertLessEqual(total, height)

    def test_falls_back_to_minimum_font_when_nothing_fits(self):
        # 어떤 크기로도 못 담으면 잘리더라도 최소 크기로 그린다(겹치는 것보다 낫다).
        font, _, _, _ = self._fit(rows=200, height_in=2.0)

        self.assertEqual(font, MIN_TABLE_FONT_SIZE)

    def test_wrapped_cell_takes_more_than_one_line(self):
        long_text = "약물치료를 받은 건선 환자 중 전신치료를 받는 환자 수와 비율" * 2
        _, row_lines, total, height = self._fit(
            rows=4, height_in=3.0, cells={(0, 0): long_text}
        )

        self.assertGreater(row_lines[0], 1, "긴 글이 든 행은 여러 줄로 잡혀야 한다")
        self.assertEqual(row_lines[1], 1)
        self.assertLessEqual(total, height)

    def test_merged_cell_spreads_its_lines_over_the_merged_rows(self):
        table = report_page(4, False, cells={(0, 0): "세로로 병합된 긴 설명 문구"},
                            merges=[(0, 0, 3, 0)]).table
        spans = _cell_spans(table)

        self.assertEqual(spans[(0, 0)], (1, 4))


class FootnoteMeasurementTests(unittest.TestCase):
    def test_counts_wrapped_lines_not_just_newlines(self):
        text = "Q6_A+B+C. Naïve, Switching, Maintain 환자별로 지난 한 달 간 처방 받은 환자 수는?" * 5

        narrow = _wrapped_line_count(text, 3.0, style.FOOTNOTE_FONT_SIZE)
        wide = _wrapped_line_count(text, 12.4, style.FOOTNOTE_FONT_SIZE)

        self.assertGreater(narrow, wide)
        self.assertGreater(wide, 1, "줄바꿈 문자가 없어도 폭을 넘으면 여러 줄로 센다")

    def test_short_footnote_stays_on_one_line(self):
        real = "Q6_A+B+C. Naïve, Switching, Maintain 환자별로 지난 한 달 간 다음의 생물학적 제제별 처방 받은 환자 수는 어떻게 되나요?"

        self.assertEqual(_wrapped_line_count(real, 12.4, style.FOOTNOTE_FONT_SIZE), 1)

    def test_empty_footnote_reserves_nothing(self):
        self.assertEqual(_wrapped_line_count("", 12.4, style.FOOTNOTE_FONT_SIZE), 0)

    def test_explicit_newlines_are_counted(self):
        self.assertEqual(_wrapped_line_count("가\n나\n다", 12.4, style.FOOTNOTE_FONT_SIZE), 3)
