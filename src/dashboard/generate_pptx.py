"""대시보드 HTML을 캡처해 PsO H-Biologics Tracker 형식의 PPTX 리포트를 생성한다.

표/텍스트/차트 모두 파워포인트에서 직접 편집 가능한 네이티브 개체(표, 텍스트박스,
차트)로 만든다. 차트 이미지(chart_png)는 PDF 전용이며, PPTX에는 Chart.js에서 그대로
읽어온 계열 데이터(chart_data)로 네이티브 차트를 그린다.
"""

import math
from functools import lru_cache
from io import BytesIO
from pathlib import Path
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE, XL_TICK_LABEL_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Pt

from src.dashboard import report_style as style
from src.dashboard.build_dashboard import _report_label
from src.dashboard.report_data import (
    ChartData,
    ReportPage,
    capture_static_pages,
    collect_report_pages,
    label_column_count,
)
from src.utils.logger import get_logger


logger = get_logger(__name__)

KOREAN_FONT = "맑은 고딕"

EMU_PER_IN = 914400
PAGE_W = Emu(int(style.PAGE_WIDTH_IN * EMU_PER_IN))
PAGE_H = Emu(int(style.PAGE_HEIGHT_IN * EMU_PER_IN))
MARGIN = Emu(int(style.MARGIN_IN * EMU_PER_IN))


# 표/차트와 각주 사이에 반드시 두는 간격. 0이면 표 아래 테두리에 각주가 붙어 겹쳐 보인다.
FOOTNOTE_GAP = Emu(int(0.14 * EMU_PER_IN))
# 파워포인트는 글자가 들어갈 높이보다 낮은 행을 그리지 못하고 행을 늘려 버린다.
# 그러면 표가 아래로 넘쳐 각주를 덮으므로, 행 높이에 맞춰 글자 크기를 먼저 줄인다.
KOREAN_LINE_HEIGHT = 1.33  # 맑은 고딕 기준 줄 높이 배수
MIN_TABLE_FONT_SIZE = 5.0
_MEASURE_FONT_SIZE = 100


def _rgb(t: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*t)


@lru_cache(maxsize=1)
def _measure_font():
    """글자 폭을 재는 용도의 폰트. 없으면 None(어림 계산으로 대체)."""
    try:
        return ImageFont.truetype(str(style.KOREAN_FONT_REGULAR_PATH), _MEASURE_FONT_SIZE)
    except OSError:
        logger.warning("⚠️ 맑은 고딕을 열지 못해 각주 줄 수를 어림 계산합니다")
        return None


def _text_width_in(text: str, font_pt: float) -> float:
    font = _measure_font()
    if font is None:
        return len(text) * font_pt / 72  # 한글 한 글자를 폰트 크기만큼으로 어림한다
    return font.getlength(text) / _MEASURE_FONT_SIZE * font_pt / 72


def _wrapped_line_count(text: str, width_in: float, font_pt: float) -> int:
    """주어진 폭에서 텍스트가 실제로 몇 줄로 접히는지 센다."""
    if not text or width_in <= 0:
        return 0
    lines = 0
    for raw_line in text.split("\n"):
        width = _text_width_in(raw_line.strip(), font_pt)
        lines += max(1, math.ceil(width / width_in))
    return lines




def _add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def _set_run(run, text, size, color, bold=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = _rgb(color)
    run.font.bold = bold
    run.font.name = KOREAN_FONT


def _add_full_bleed_picture(slide, png_bytes: bytes, crop_to_fill: bool = False):
    img = Image.open(BytesIO(png_bytes))
    iw, ih = img.size
    scale_fn = max if crop_to_fill else min
    scale = scale_fn(PAGE_W / iw, PAGE_H / ih)
    w, h = int(iw * scale), int(ih * scale)
    x, y = (PAGE_W - w) // 2, (PAGE_H - h) // 2
    slide.shapes.add_picture(BytesIO(png_bytes), x, y, width=w, height=h)


def _add_static_slide(prs: Presentation, png_bytes: bytes, crop_to_fill: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_full_bleed_picture(slide, png_bytes, crop_to_fill=crop_to_fill)


def _hex(color: str) -> str:
    return color.lstrip("#").upper() or "808080"


def _add_native_chart(slide, chart_data: ChartData, left, top, width, height):
    """Chart.js 데이터를 그대로 옮겨 파워포인트 네이티브 차트(막대, 필요하면 보조축 꺾은선)로 그린다."""
    bar_series = [s for s in chart_data.series if s.kind != "line"]
    line_series = [s for s in chart_data.series if s.kind == "line"]

    cdata = CategoryChartData()
    cdata.categories = chart_data.categories
    for s in bar_series:
        cdata.add_series(s.label, [v if v is not None else 0 for v in s.values])

    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, cdata)
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(7)
    chart.legend.font.name = KOREAN_FONT

    plot = chart.plots[0]
    plot.gap_width = 40
    plot.overlap = 100
    for series, s in zip(plot.series, bar_series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor.from_string(_hex(s.color))
        series.format.line.fill.background()

    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(7)
        axis.tick_labels.font.name = KOREAN_FONT

    if line_series:
        try:
            _add_secondary_line_chart(slide, line_series, chart_data.categories, left, top, width, height)
        except Exception:
            logger.exception("⚠️ 보조축 꺾은선 계열 추가 실패 — 막대 차트만 유지합니다")

    return graphic_frame


def _add_secondary_line_chart(slide, line_series: list, categories: list[str], left, top, width, height):
    """편집 가능한 투명 꺾은선 차트를 막대 차트 위에 배치한다.

    콤보 차트를 만들기 위한 OOXML 직접 삽입은 PowerPoint 복구를 유발할 수 있으므로
    python-pptx가 보장하는 정식 차트 객체 두 개를 사용한다.
    """
    line_data = CategoryChartData()
    line_data.categories = categories
    for series in line_series:
        line_data.add_series(series.label, [value if value is not None else 0 for value in series.values])

    frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, line_data)
    line_chart = frame.chart
    line_chart.has_legend = False
    _make_chart_background_transparent(line_chart)
    line_chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    line_chart.category_axis.format.line.fill.background()
    line_chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
    line_chart.value_axis.format.line.fill.background()
    line_chart.value_axis.has_major_gridlines = False
    line_chart.value_axis.minimum_scale = 0
    line_chart.value_axis.maximum_scale = 100

    for native_series, source in zip(line_chart.series, line_series):
        color = RGBColor.from_string(_hex(source.color))
        native_series.format.line.color.rgb = color
        native_series.format.line.width = Pt(1.5)
        native_series.marker.style = XL_MARKER_STYLE.CIRCLE
        native_series.marker.size = 5
        native_series.marker.format.fill.solid()
        native_series.marker.format.fill.fore_color.rgb = color
        native_series.marker.format.line.color.rgb = color

    return frame


def _make_chart_background_transparent(chart) -> None:
    """차트 영역과 플롯 영역에 표준 noFill을 적용한다."""
    chart_space = chart._chartSpace
    chart_element = chart_space.chart
    sp_pr = parse_xml(
        f'<c:spPr {nsdecls("c", "a")}><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'
    )
    chart_element.addnext(sp_pr)

    plot_area = chart_element.plotArea
    plot_sp_pr = parse_xml(
        f'<c:spPr {nsdecls("c", "a")}><a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>'
    )
    plot_area.append(plot_sp_pr)


def _table_height_for(page: ReportPage, available: Emu, has_chart: bool) -> Emu:
    """행 수에 맞춰 표 높이를 정하되 차트와 함께 가용 영역을 넘지 않게 한다."""
    n_rows = page.table.n_rows if page.table else 0
    if not n_rows:
        return Emu(0)
    preferred = Emu(int(n_rows * 0.22 * EMU_PER_IN))
    if has_chart:
        # 차트가 있는 장표는 표가 콘텐츠 영역의 절반 이상을 선점하지 않게 한다.
        limit = Emu(int(available * 0.48))
    else:
        limit = available
    return min(preferred, limit)


def _cell_spans(table_data) -> dict[tuple[int, int], tuple[int, int]]:
    """병합된 칸의 (colspan, rowspan)을 모은다. 병합 칸은 그만큼 넓고 높아 글이 덜 접힌다."""
    spans: dict[tuple[int, int], tuple[int, int]] = {}
    for (r1, c1, r2, c2) in table_data.merges:
        spans[(r1, c1)] = (c2 - c1 + 1, r2 - r1 + 1)
    return spans


def _row_line_counts(table_data, col_widths, spans, font_pt: float) -> list[int]:
    """글자 크기를 정했을 때 각 행이 몇 줄을 차지하는지 센다."""
    h_margin_in = 2 * 18000 / EMU_PER_IN
    lines_per_row = [1] * table_data.n_rows
    for (r, c), text in table_data.cells.items():
        content = text.replace("\n", " ").strip()
        if not content:
            continue
        colspan, rowspan = spans.get((r, c), (1, 1))
        width_in = sum(col_widths[c : c + colspan]) / EMU_PER_IN - h_margin_in
        if width_in <= 0:
            continue
        lines = max(1, math.ceil(_text_width_in(content, font_pt) / width_in))
        # 세로로 병합된 칸은 여러 행에 걸쳐 있으므로 행마다 필요한 줄은 그만큼 적다.
        per_row = max(1, math.ceil(lines / rowspan))
        for offset in range(rowspan):
            if r + offset < len(lines_per_row):
                lines_per_row[r + offset] = max(lines_per_row[r + offset], per_row)
    return lines_per_row


def _fit_table_font(table_data, col_widths, spans, height: Emu):
    """표 전체가 height 안에 들어가는 글자 크기·여백·행별 줄 수를 찾는다."""
    for step in range(0, 21):  # 7.5pt에서 0.125pt씩 낮추며 맞는 크기를 찾는다
        font_pt = style.TABLE_FONT_SIZE - step * 0.125
        if font_pt < MIN_TABLE_FONT_SIZE:
            break
        for v_margin in (Emu(4000), Emu(0)):
            row_lines = _row_line_counts(table_data, col_widths, spans, font_pt)
            line_h = font_pt * KOREAN_LINE_HEIGHT / 72 * EMU_PER_IN
            total = sum(lines * line_h + 2 * int(v_margin) for lines in row_lines)
            if total <= int(height):
                return font_pt, v_margin, row_lines

    # 아무리 줄여도 안 들어가면 최소 크기로 그리고 알린다(표가 잘리는 편이 겹치는 것보다 낫다).
    row_lines = _row_line_counts(table_data, col_widths, spans, MIN_TABLE_FONT_SIZE)
    logger.warning(
        "⚠️ 표가 너무 길어 %.1fpt로도 영역에 맞추지 못했습니다 (%d행)",
        MIN_TABLE_FONT_SIZE, table_data.n_rows,
    )
    return MIN_TABLE_FONT_SIZE, Emu(0), row_lines


def _build_table(slide, page: ReportPage, top: Emu, height: Emu) -> Emu:
    """표를 그리고 실제로 차지한 높이를 반환한다."""
    table_data = page.table
    label_cols = label_column_count(table_data)
    n_rows, n_cols = table_data.n_rows, table_data.n_cols

    total_width = PAGE_W - 2 * MARGIN
    label_frac = 0.17 if label_cols == 1 else 0.26
    label_total_w = int(total_width * label_frac)
    data_col_w = (total_width - label_total_w) // max(1, n_cols - label_cols)
    label_col_w = label_total_w // label_cols

    col_widths = [label_col_w if c < label_cols else data_col_w for c in range(n_cols)]
    spans = _cell_spans(table_data)

    # 파워포인트는 글자가 들어갈 높이보다 낮은 행을 그리지 못하고 행을 늘려 버린다.
    # 그러면 표가 아래로 넘쳐 각주를 덮으므로, 긴 글이 접히는 줄 수까지 계산해
    # 표 전체가 주어진 높이 안에 확실히 들어가는 글자 크기를 먼저 찾는다.
    font_size, v_margin, row_lines = _fit_table_font(table_data, col_widths, spans, height)

    line_h = Emu(int(font_size * KOREAN_LINE_HEIGHT / 72 * EMU_PER_IN))
    row_heights = [Emu(int(lines * line_h + 2 * v_margin)) for lines in row_lines]
    # 남는 공간은 모든 행에 고르게 나눠 표가 영역을 자연스럽게 채우게 한다.
    slack = int(height) - sum(int(h) for h in row_heights)
    if slack > 0:
        share = slack // n_rows
        row_heights = [Emu(int(h) + share) for h in row_heights]
    table_height = Emu(sum(int(h) for h in row_heights))

    gframe = slide.shapes.add_table(n_rows, n_cols, MARGIN, top, total_width, table_height)
    tbl = gframe.table
    for c in range(n_cols):
        tbl.columns[c].width = col_widths[c]
    for r in range(n_rows):
        tbl.rows[r].height = row_heights[r]

    for (r, c), text in table_data.cells.items():
        cell = tbl.cell(r, c)
        cell.text = text.replace("\n", " ")
        cell.margin_left = cell.margin_right = Emu(18000)
        cell.margin_top = cell.margin_bottom = v_margin
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        is_header = r < table_data.header_rows
        is_bold_row = (r - table_data.header_rows) in table_data.bold_rows
        is_label = c < label_cols
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT if is_label else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(font_size)
                run.font.name = KOREAN_FONT
                if is_header:
                    run.font.bold = True
                    run.font.color.rgb = _rgb(style.TABLE_HEADER_TEXT)
                elif is_bold_row:
                    run.font.bold = True
                    run.font.color.rgb = _rgb(style.TABLE_BOLD_ROW_LABEL_TEXT if is_label else style.TABLE_BODY_TEXT)
                else:
                    run.font.color.rgb = _rgb(style.TABLE_BODY_TEXT)
        cell.fill.solid()
        if is_header:
            bg = style.TABLE_HEADER_TOP_BG if r == 0 else style.TABLE_HEADER_PER_BG
        elif is_bold_row:
            bg = style.TABLE_BOLD_ROW_LABEL_BG if is_label else style.TABLE_BOLD_ROW_BG
        else:
            bg = (0xFF, 0xFF, 0xFF)
        cell.fill.fore_color.rgb = _rgb(bg)

    for (r1, c1, r2, c2) in table_data.merges:
        if (r1, c1) != (r2, c2):
            tbl.cell(r1, c1).merge(tbl.cell(r2, c2))

    return table_height


def _add_data_slide(prs: Presentation, page: ReportPage):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    x = MARGIN
    y = MARGIN

    # 강조 바
    bar = slide.shapes.add_shape(1, x, y, Emu(int(0.045 * EMU_PER_IN)), Emu(int(0.22 * EMU_PER_IN)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(style.ACCENT_BAR)
    bar.line.fill.background()

    # 제목
    _, tf = _add_textbox(slide, x + Emu(int(0.12 * EMU_PER_IN)), y, PAGE_W - 2 * MARGIN, Emu(int(0.3 * EMU_PER_IN)))
    _set_run(tf.paragraphs[0].add_run(), page.title, style.TITLE_FONT_SIZE, style.TITLE_TEXT, bold=True)
    y += Emu(int(0.34 * EMU_PER_IN))

    # 인사이트
    if page.insight_lines:
        insight_height = Emu(int((0.02 + 0.19 * len(page.insight_lines)) * EMU_PER_IN))
        _, tf = _add_textbox(slide, x + Emu(int(0.12 * EMU_PER_IN)), y, PAGE_W - 2 * MARGIN,
                                    insight_height)
        for i, line in enumerate(page.insight_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_run(p.add_run(), f"▶ {line}", style.INSIGHT_FONT_SIZE, style.INSIGHT_TEXT)
        y += insight_height

    # Unit note
    if page.unit_note:
        _, tf = _add_textbox(slide, x, y, PAGE_W - 2 * MARGIN, Emu(int(0.16 * EMU_PER_IN)))
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _set_run(tf.paragraphs[0].add_run(), page.unit_note, style.UNIT_NOTE_FONT_SIZE, style.UNIT_NOTE_TEXT)
        y += Emu(int(0.2 * EMU_PER_IN))

    y += Emu(int(0.04 * EMU_PER_IN))

    footnote_lines = (page.footnote or "").split("\n") if page.footnote else []
    # 각주는 슬라이드 폭에 맞춰 접히므로, 줄바꿈 문자 수가 아니라 실제로 그려질 줄 수로 자리를 잡는다.
    text_width_in = (PAGE_W - 2 * MARGIN) / EMU_PER_IN
    visual_lines = _wrapped_line_count(page.footnote or "", text_width_in, style.FOOTNOTE_FONT_SIZE)
    bottom_reserved = Emu(int((0.10 + 0.13 * visual_lines) * EMU_PER_IN))
    footnote_top = PAGE_H - MARGIN - bottom_reserved
    # 표/차트가 각주에 닿지 않도록 사이에 간격을 둔다.
    content_bottom = footnote_top - (FOOTNOTE_GAP if page.footnote else Emu(0))
    content_height = content_bottom - y

    gap = Emu(int(0.06 * EMU_PER_IN))
    table_h = _table_height_for(page, content_height, bool(page.chart_png))
    chart_h = Emu(0)
    if page.chart_png:
        img = Image.open(BytesIO(page.chart_png))
        iw, ih = img.size
        avail_w = PAGE_W - 2 * MARGIN
        chart_available = max(Emu(0), content_height - table_h - gap)
        chart_h = min(chart_available, Emu(int(avail_w * ih / iw)))
        chart_w = Emu(int(chart_h * iw / ih))
        if chart_w > avail_w:
            chart_w = avail_w
            chart_h = Emu(int(chart_w * ih / iw))
        chart_x = x + (avail_w - chart_w) // 2

        if page.chart_data and page.chart_data.series:
            try:
                _add_native_chart(slide, page.chart_data, chart_x, y, chart_w, chart_h)
            except Exception:
                logger.exception("⚠️ 네이티브 차트 생성 실패 — 이미지로 대체합니다 (%s)", page.section_id)
                slide.shapes.add_picture(BytesIO(page.chart_png), chart_x, y, width=chart_w, height=chart_h)
        else:
            slide.shapes.add_picture(BytesIO(page.chart_png), chart_x, y, width=chart_w, height=chart_h)
        y += chart_h + gap

    if page.table is not None:
        remaining = max(Emu(0), content_bottom - y)
        _build_table(slide, page, y, min(table_h, remaining))

    # 각주
    if page.footnote:
        _, tf = _add_textbox(slide, x, footnote_top, PAGE_W - 2 * MARGIN, bottom_reserved)
        for i, line in enumerate(footnote_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_run(p.add_run(), line.strip(), style.FOOTNOTE_FONT_SIZE, style.FOOTNOTE_TEXT)

    # 페이지 번호
    _, tf = _add_textbox(slide, PAGE_W - MARGIN - Emu(int(0.6 * EMU_PER_IN)), PAGE_H - MARGIN,
                                 Emu(int(0.6 * EMU_PER_IN)), Emu(int(0.18 * EMU_PER_IN)))
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _set_run(tf.paragraphs[0].add_run(), page.page_num, style.PAGE_NUM_FONT_SIZE, style.PAGE_NUM_TEXT)


def generate_pptx(output_path: str, dashboard_html: str, period: str) -> Path:
    """대시보드 HTML을 캡처해 PPTX 리포트를 생성한다.

    Args:
        output_path: 생성할 PPTX 파일 경로
        dashboard_html: build_dashboard()가 방금 만든 대시보드 HTML 경로
        period: 보고 차수 (예: "26년 6차")
    """
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    report_date_label = _report_label(period)
    logger.info("📊 PPTX 리포트용 장표 캡처 시작 (%s)", period)
    static_pages = capture_static_pages(dashboard_html, report_date_label)
    data_pages = collect_report_pages(dashboard_html)
    logger.info("📊 장표 %d개 캡처 완료, PPTX 조립 시작", len(data_pages))

    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    _add_static_slide(prs, static_pages["cover"], crop_to_fill=True)
    _add_static_slide(prs, static_pages["toc_survey"])
    _add_static_slide(prs, static_pages["overview"])
    _add_static_slide(prs, static_pages["toc_result"])
    for page in data_pages:
        _add_data_slide(prs, page)

    prs.save(str(output))
    logger.info("✅ PPTX 리포트 저장 완료: %s", output)
    return output
